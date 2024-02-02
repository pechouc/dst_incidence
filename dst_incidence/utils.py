import os

import time
# from datetime import date

import re

import numpy as np
import pandas as pd

import matplotlib.pyplot as plt

from pptx.util import Inches


def clean_date_string(string):

    string = string.split('(')[1].strip().strip(')').replace('.', '')

    day, fr_month, year = string.split(' ')

    month = {
        'jan': 'jan',
        'fév': 'feb',
        'mar': 'mar',
        'avr': 'apr',
        'mai': 'may',
        'juin': 'jun',
        'juil': 'jul',
        'août': 'aug',
        'sept': 'sep',
        'oct': 'oct',
        'nov': 'nov',
        'déc': 'dec',
    }.get(fr_month, fr_month)

    string = ' '.join([day, month, year])

    return string


def draw_random_sku_country(path_to_price_data):

    path_to_price_history = os.path.join(path_to_price_data, 'price_history.csv')
    path_to_price_history_summary = os.path.join(path_to_price_data, 'price_history_summary.csv')

    sku_country = pd.read_csv(path_to_price_history, usecols=["sku", "country"])

    random_draw = sku_country.sample(1, random_state=int(time.time())).reset_index(drop=True)
    random_sku, random_country = random_draw.loc[0, 'sku'], random_draw.loc[0, 'country']

    print("SKU:", random_sku, "---", "Country:", random_country)

    price_history = pd.read_csv(
        path_to_price_history,
        dtype={"trusted": bool, "price": float, "date": str, "sku": str, "country": str, "seller": str},
        parse_dates=['date']
    )
    price_history = price_history[
        np.logical_and(
            price_history['sku'] == random_sku,
            price_history['country'] == random_country
        )
    ].copy()

    price_history['date'] = pd.to_datetime(price_history['date'])

    new_price_history = price_history[price_history['seller'] == 'new'].copy()
    amazon_price_history = price_history[price_history['seller'] == 'amazon'].copy()

    del price_history

    new_price_history['transition'] = np.logical_and(
        new_price_history['trusted'] != new_price_history['trusted'].shift(1),
        ~new_price_history['trusted'].shift(1).isnull()
    )

    new_price_history['transition#'] = new_price_history['transition'].cumsum()

    new_trusted_dfs = []
    new_not_trusted_dfs = []

    if not new_price_history.empty:

        for i in range(new_price_history['transition#'].max() + 1):
            extract = new_price_history[new_price_history['transition#'] == i].copy()

            if extract['trusted'].unique()[0]:
                new_trusted_dfs.append(extract)

            else:
                new_not_trusted_dfs.append(extract)

    amazon_price_history['transition'] = np.logical_and(
        amazon_price_history['trusted'] != amazon_price_history['trusted'].shift(1),
        ~amazon_price_history['trusted'].shift(1).isnull()
    )

    amazon_price_history['transition#'] = amazon_price_history['transition'].cumsum()

    amazon_trusted_dfs = []
    amazon_not_trusted_dfs = []

    price_history_summary = pd.read_csv(path_to_price_history_summary)
    price_history_summary = price_history_summary[
        np.logical_and(
            price_history_summary['sku'] == random_sku,
            price_history_summary['country'] == random_country
        )
    ].copy()

    amazon_price_history_summary = price_history_summary[
        price_history_summary["seller"] == "amazon"
    ].reset_index(drop=True)
    new_price_history_summary = price_history_summary[
        price_history_summary["seller"] == "new"
    ].reset_index(drop=True)

    max_prices = {
        'amazon': amazon_price_history_summary.loc[0, 'max_price'],
        'new': new_price_history_summary.loc[0, 'max_price']
    }
    min_prices = {
        'amazon': amazon_price_history_summary.loc[0, 'min_price'],
        'new': new_price_history_summary.loc[0, 'min_price']
    }

    if not amazon_price_history.empty:

        for i in range(amazon_price_history['transition#'].max() + 1):
            extract = amazon_price_history[amazon_price_history['transition#'] == i].copy()

            if bool(extract['trusted'].unique()[0]):
                amazon_trusted_dfs.append(extract)

            else:
                amazon_not_trusted_dfs.append(extract)

    return (
        random_sku,
        random_country,
        new_trusted_dfs.copy(),
        new_not_trusted_dfs.copy(),
        amazon_trusted_dfs.copy(),
        amazon_not_trusted_dfs.copy(),
        max_prices.copy(),
        min_prices.copy()
    )


def get_dates_from_text_extract(text_extract):

    # Getting the first month
    months = np.array(
        ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    )

    months_idx = []

    for month in months:
        months_idx.append(
            text_extract.find(month) if text_extract.find(month) > -1 else np.nan
        )

    months_idx = np.array(months_idx)

    first_month = months[np.where(months_idx == np.nanmin(months_idx))][0]

    # Getting the last month
    months_idx_right = []

    for month in months:
        months_idx_right.append(
            text_extract.rfind(month) if text_extract.rfind(month) > -1 else np.nan
        )

    months_idx_right = np.array(months_idx_right)

    last_month = months[np.where(months_idx_right == np.nanmax(months_idx_right))][0]

    # Getting first and last years
    pattern = r'\b(\d{2})[°.]?[ \t](?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)'

    # potential_first_year_matches = [
    #     match for match in re.finditer(r'[^€\d+,](\d\d)(?![,;]\d{2})', text_extract)
    # ]
    potential_first_year_matches = [
        match for match in re.finditer(pattern, text_extract)
    ]
    potential_first_year_matches = [
        match for match in potential_first_year_matches if int(
            ''.join(filter(str.isdigit, match.group(0)))
        ) <= 20
    ]
    potential_first_year_matches = [
        match for match in potential_first_year_matches if int(
            ''.join(filter(str.isdigit, match.group(0)))
        ) >= 7
    ]
    potential_first_years = [
        int(''.join(filter(str.isdigit, match.group(0)))) for match in potential_first_year_matches
    ]
    first_year = 2000 + int(potential_first_years[0])

    # last_year = 2000 + int(re.findall(r'[^€\d+,](\d\d)(?![,;]\d{2})', text_extract)[-1])
    last_year = 2000 + int(re.findall(pattern, text_extract)[-1])
    last_year_matches = [match for match in re.finditer(r'[^€\d+,](\d\d)(?![,;]\d{2})', text_extract)]
    # last_year_matches = [match for match in re.finditer(pattern, text_extract)]
    last_year_match = last_year_matches[-1]
    last_year_start = last_year_match.start()

    # Correcting first and last months if needed
    if potential_first_year_matches[0].start() < text_extract.find(first_month):
        first_month = 'Jan'
    else:
        first_year = first_year - 1
    if last_year_start > text_extract.rfind(last_month):
        last_month = 'Jan'

    return first_month, first_year, last_month, last_year


# Function to add a slide with side-by-side image and plot
def add_slide_with_image_and_plot(prs, image_path, matplotlib_plot):
    slide_layout = prs.slide_layouts[5]  # Use a content with caption layout (could be different based on your preference)

    # Remove the title placeholder from the slide layout
    for shape in slide_layout.placeholders:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            shape.text = ''
            shape.element.clear()  # Clear the existing elements

    # Add a slide
    slide = prs.slides.add_slide(slide_layout)

    # Define the position and size of the image and plot
    left_inch = Inches(0.5)   # Adjust as needed
    top_inch = Inches(0.5)    # Adjust as needed
    width_inch = Inches(7)    # Adjust as needed
    height_inch = Inches(7)   # Adjust as needed

    # Add the image
    slide.shapes.add_picture(image_path, left_inch, top_inch, width_inch, height_inch)

    # Save the matplotlib plot to a temporary file
    temp_plot_file = "temp_plot.png"
    matplotlib_plot.savefig(temp_plot_file, bbox_inches='tight')

    # Add the matplotlib plot to the slide
    slide.shapes.add_picture(temp_plot_file, left_inch + width_inch + Inches(0.5), top_inch, width_inch, height_inch)

    # Clean up the temporary plot file
    plt.close()
    import os
    os.remove(temp_plot_file)


def create_own_graph_for_random_check(
        seller,
        new_trusted_dfs,
        new_not_trusted_dfs,
        amazon_trusted_dfs,
        amazon_not_trusted_dfs,
        max_prices,
        min_prices
    ):
    if seller == "amazon":
        trusted_dfs = amazon_trusted_dfs.copy()
        not_trusted_dfs = amazon_not_trusted_dfs.copy()
    else:
        trusted_dfs = new_trusted_dfs.copy()
        not_trusted_dfs = new_not_trusted_dfs.copy()

    fig, ax = plt.subplots(figsize=(15, 11))

    min_date_trusted = pd.to_datetime('2026-12-04')
    max_date_trusted = pd.to_datetime('2001-12-04')
    min_date_not_trusted = pd.to_datetime('2026-12-04')
    max_date_not_trusted = pd.to_datetime('2001-12-04')

    if len(trusted_dfs) + len(not_trusted_dfs) > 0:

        for i, trusted_df in enumerate(trusted_dfs):
            if i == 0:
                min_date_trusted = trusted_df['date'].min()
            if i == len(trusted_dfs) - 1:
                max_date_trusted = trusted_df['date'].max()

            ax.plot(
                trusted_df['date'],
                trusted_df['price'],
                color='darkblue' if seller == "new" else "green"
            )

        for i, not_trusted_df in enumerate(not_trusted_dfs):
            if i == 0:
                min_date_not_trusted = not_trusted_df['date'].min()
            if i == len(not_trusted_dfs) - 1:
                max_date_not_trusted = not_trusted_df['date'].max()

            ax.scatter(
                not_trusted_df['date'],
                not_trusted_df['price'],
                marker='s',
                s=1,
                color='darkred'
            )

        ax.hlines(
            (max_prices[seller], min_prices[seller]),
            xmin=min(min_date_trusted, min_date_not_trusted),
            xmax=max(max_date_trusted, max_date_not_trusted),
            colors=('red', 'green'),
            linestyles='dashed',
        )

    else:

        ax.plot([1, 2], [1, 2], color="orange")
        ax.set_title("Skipped due to lack of data")

    return fig, ax
