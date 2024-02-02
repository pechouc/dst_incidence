import os

from dst_incidence.utils import (
    draw_random_sku_country,
    add_slide_with_image_and_plot,
    create_own_graph_for_random_check,
)

import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches


if __name__ == "__main__":

    # Create a presentation object
    prs = Presentation()

    # Set the slide width and height (wide layout)
    prs.slide_width = Inches(16)  # Width of the slide (in inches)
    prs.slide_height = Inches(9)   # Height of the slide (in inches)

    path_to_dir = os.path.dirname(os.path.abspath(__file__))
    path_to_price_data = os.path.join(path_to_dir, 'data', 'prices')

    path_to_data_raw = "/Users/Paul-Emmanuel/Desktop/PhD/3_DST_incidence/data_raw"

    for _ in range(20):

        (
            random_sku,
            random_country,
            new_trusted_dfs,
            new_not_trusted_dfs,
            amazon_trusted_dfs,
            amazon_not_trusted_dfs,
            max_prices,
            min_prices
        ) = draw_random_sku_country(path_to_price_data=path_to_price_data)

        # Add a section title slide
        slide_layout = prs.slide_layouts[0]  # Use the layout suitable for section titles
        slide = prs.slides.add_slide(slide_layout)

        # Set the title and subtitle (change as per your requirement)
        title = slide.shapes.title
        title.text = f"{random_sku} - {random_country}"

        for seller in ['new', 'amazon']:

            paths_to_chart = os.path.join(
                path_to_data_raw, f"charts_{seller}_{random_country}", f"{random_sku}.png"
            )

            # --- Graphs based on our data

            fig, ax = create_own_graph_for_random_check(
                seller=seller,
                new_trusted_dfs=new_trusted_dfs,
                new_not_trusted_dfs=new_not_trusted_dfs,
                amazon_trusted_dfs=amazon_trusted_dfs,
                amazon_not_trusted_dfs=amazon_not_trusted_dfs,
                max_prices=max_prices,
                min_prices=min_prices
            )

            add_slide_with_image_and_plot(prs, paths_to_chart, fig)

    prs.save(
        os.path.join(
            os.path.dirname(path_to_dir),
            "tests",
            "presentation_for_random_checks.pptx"
        )
    )
